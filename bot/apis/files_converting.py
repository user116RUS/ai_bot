#from bot import logger

import PyPDF2
import pdfplumber
from docx import Document
import subprocess
from pptx import Presentation
import openpyxl
import xlrd
from bs4 import BeautifulSoup


class DocumentConverter:
    def __init__(self):
        pass

    def pdf_to_text(self, file_path):
        """Convert PDF to text."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                with pdfplumber.open(file) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
        except Exception as e:
            #logger.error(f"Error reading PDF: {e}")
            print(f"Error reading PDF: {e}")
        return text

    def docx_to_text(self, file_path):
        """Convert DOCX to text."""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            #logger.error(f"Error reading DOCX: {e}")
            print(f"Error reading DOCX: {e}")
        return text

    def pptx_to_text(self, file_path):
        """Convert PPTX to text."""
        text = ""
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            #logger.error(f"Error reading PPTX: {e}")
            print(f"Error reading PPTX: {e}")
        return text

    def xlsx_to_text(self, file_path):
        """Convert XLSX to text."""
        text = ""
        try:
            wb = openpyxl.load_workbook(file_path)
            for sheet in wb.sheetnames:
                worksheet = wb[sheet]
                for row in worksheet.iter_rows(values_only=True):
                    text += "\t".join(str(cell) for cell in row) + "\n"
        except Exception as e:
            #logger.error(f"Error reading XLSX: {e}")
            print(f"Error reading XLSX: {e}")
        return text

    def xls_to_text(self, file_path):
        """Convert XLS (older version) to text."""
        text = ""
        try:
            wb = xlrd.open_workbook(file_path)
            for sheet in wb.sheets():
                for row in range(sheet.nrows):
                    text += "\t".join(str(cell) for cell in sheet.row(row)) + "\n"
        except Exception as e:
            #logger.error(f"Error reading XLS: {e}")
            print(f"Error reading XLS: {e}")
        return text

    def html_to_text(self, file_path):
        """Convert HTML file to text."""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                text = soup.get_text()
        except Exception as e:
            #logger.error(f"Error reading HTML: {e}")
            print(f"Error reading HTML: {e}")
        return text

    def txt_to_text(self, file_path):
        """Convert plain text file to text."""
        text = ""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        except Exception as e:
            #logger.error(f"Error reading TXT: {e}")
            print(f"Error reading TXT: {e}")
        return text

    def convert(self, file_path):
        """Automatically detect file type and convert to text."""
        if str(file_path).endswith('.pdf'):
            return self.pdf_to_text(file_path)
        elif file_path.endswith('.docx'):
            return self.docx_to_text(file_path)
        elif file_path.endswith('.pptx'):
            return self.pptx_to_text(file_path)
        elif file_path.endswith('.xlsx'):
            return self.xlsx_to_text(file_path)
        elif file_path.endswith('.xls'):
            return self.xls_to_text(file_path)
        elif file_path.endswith('.html'):
            return self.html_to_text(file_path)
        elif file_path.endswith('.txt'):
            return self.txt_to_text(file_path)
        else:
            raise ValueError("Unsupported file type")
